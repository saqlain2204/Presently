import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from most_similar_image import find_best_matching_image
from generate_image import generate_image

def parse_markdown(text):
    lines = text.splitlines()
    title = None
    slides = []
    current_slide = None
    
    for line in lines:
        line = line.strip()
        if line.startswith('## ') and title is None:
            title = line[3:].strip()
        elif line.startswith('# '):
            if current_slide:
                slides.append(current_slide)
            current_slide = {'title': line[2:].strip(), 'points': []}
        elif line.startswith('- ') and current_slide:
            current_slide['points'].append(line[2:].strip())
    
    if current_slide:
        slides.append(current_slide)
    
    return title, slides

def add_title_slide(prs, title_text):
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255) 

    left = Inches(0)
    top = Inches(3)
    width = prs.slide_width
    height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    font = run.font
    font.size = Pt(48)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, points, image_path):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    font = run.font
    font.size = Pt(32)
    font.bold = True
    font.color.rgb = RGBColor(0, 0, 0)
    
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(4.3)
    height = Inches(4.5)
    body_box = slide.shapes.add_textbox(left, top, width, height)
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.bullet = True
        p.space_after = Pt(10)
    
    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(4.5)
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)
        except Exception as e:
            print(f"Could not add image '{image_path}': {e}")
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.text = "Image\nNot Found"
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.text = "Image Placeholder"
    
    return slide

def add_thank_you_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    left = Inches(0)
    top = Inches(3)
    width = prs.slide_width
    height = Inches(2)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You"
    font = run.font
    font.size = Pt(40)
    font.bold = True
    font.color.rgb = RGBColor(0, 128, 0) 
    p.alignment = PP_ALIGN.CENTER
    return slide



def markdown_to_ppt(workspace_dir, output_file="presentation.ppt"):
    md_path = os.path.join(workspace_dir, "temp", "presentation.md")
    images_dir = os.path.join(workspace_dir, "temp", "images")
    
    with open(md_path, "r", encoding="utf-8") as f:
        markdown_text = f.read()
    
    title, slides = parse_markdown(markdown_text)
    prs = Presentation()
    
    add_title_slide(prs, title if title else "Presentation")
    
    for slide_data in slides:
        check_text = " ".join(slide_data['points'])
        images_dir = os.path.join(workspace_dir, "temp", "images")

        image_path = find_best_matching_image(images_dir, check_text)

        if not image_path:
            generated_image_path = os.path.join(
                images_dir, f"{slide_data['title'].replace(' ', '_')}.png"
            )
            image_path = generate_image(check_text, generated_image_path)

        add_content_slide(prs, slide_data['title'], slide_data['points'], image_path)

        if image_path and os.path.exists(image_path):
            os.remove(image_path)
    
    add_thank_you_slide(prs)
    
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")
